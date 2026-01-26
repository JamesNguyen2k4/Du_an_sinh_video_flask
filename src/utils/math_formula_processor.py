import re
import unicodedata
import xml.etree.ElementTree as ET
from typing import List, Dict, Tuple, Optional
import logging

# Thiết lập logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    MSO_SHAPE_TYPE = None


class MathFormulaProcessor:
    """
    Xử lý các ký tự đặc biệt và công thức toán học từ PowerPoint
    """
    def _mathiness_score(self, s: str) -> int:
        """Điểm 'tính toán học' theo số toán tử/ký hiệu đặc trưng."""
        if not s:
            return 0
        ops = r"[+\-*/=^×÷⋅·]|∑|∏|√|∫|≈|≠|≤|≥|⊂|⊃|∈|∉"
        funcs = r"\b(sin|cos|tan|log|ln|lim|exp)\b"
        score = len(re.findall(ops, s))
        score += 2 * len(re.findall(funcs, s, flags=re.IGNORECASE))
        # có superscript/subscript unicode?
        if re.search(r"[⁰¹²³⁴⁵⁶⁷⁸⁹₀₁₂₃₄₅₆₇₈₉]", s):
            score += 2
        return score

    def _is_math_line(self, s: str) -> bool:
        """Heuristic: coi là toán nếu có >=2 toán tử/ký hiệu, hoặc có dấu ^, ∑, ∫, √…"""
        if not s:
            return False
        if re.search(r"[∑∏√∫^]", s):
            return True
        return self._mathiness_score(s) >= 2

    def __init__(self):
        # Bảng tra cứu ký tự đặc biệt sang tiếng Việt
        self.special_char_map = {
            # Số mũ
            '²': ' mũ hai',
            '³': ' mũ ba',
            '¹': ' mũ một',
            '⁴': ' mũ bốn',
            '⁵': ' mũ năm',
            '⁶': ' mũ sáu',
            '⁷': ' mũ bảy',
            '⁸': ' mũ tám',
            '⁹': ' mũ chín',
            '⁰': ' mũ không',

            # Chỉ số dưới
            '₁': ' chỉ số một',
            '₂': ' chỉ số hai',
            '₃': ' chỉ số ba',
            '₄': ' chỉ số bốn',
            '₅': ' chỉ số năm',
            '₆': ' chỉ số sáu',
            '₇': ' chỉ số bảy',
            '₈': ' chỉ số tám',
            '₉': ' chỉ số chín',
            '₀': ' chỉ số không',

            # Ký tự Hy Lạp
            'α': ' alpha',
            'β': ' beta',
            'γ': ' gamma',
            'δ': ' delta',
            'ε': ' epsilon',
            'ζ': ' zeta',
            'η': ' eta',
            'θ': ' theta',
            'ι': ' iota',
            'κ': ' kappa',
            'λ': ' lambda',
            'μ': ' mu',
            'ν': ' nu',
            'ξ': ' xi',
            'ο': ' omicron',
            'π': ' pi',
            'ρ': ' rho',
            'σ': ' sigma',
            'τ': ' tau',
            'υ': ' upsilon',
            'φ': ' phi',
            'χ': ' chi',
            'ψ': ' psi',
            'ω': ' omega',

            # Chữ hoa Hy Lạp
            'Α': ' Alpha',
            'Β': ' Beta',
            'Γ': ' Gamma',
            'Δ': ' Delta',
            'Ε': ' Epsilon',
            'Ζ': ' Zeta',
            'Η': ' Eta',
            'Θ': ' Theta',
            'Ι': ' Iota',
            'Κ': ' Kappa',
            'Λ': ' Lambda',
            'Μ': ' Mu',
            'Ν': ' Nu',
            'Ξ': ' Xi',
            'Ο': ' Omicron',
            'Π': ' Pi',
            'Ρ': ' Rho',
            'Σ': ' Sigma',
            'Τ': ' Tau',
            'Υ': ' Upsilon',
            'Φ': ' Phi',
            'Χ': ' Chi',
            'Ψ': ' Psi',
            'Ω': ' Omega',

            # Ký tự toán học
            '±': ' cộng trừ',
            '∓': ' trừ cộng',
            '×': ' nhân',
            '÷': ' chia',
            '⋅': ' nhân',
            '∗': ' nhân',
            '√': ' căn bậc hai',
            '∛': ' căn bậc ba',
            '∜': ' căn bậc bốn',
            '∞': ' vô cùng',
            '≈': ' xấp xỉ',
            '≠': ' khác',
            '≤': ' nhỏ hơn hoặc bằng',
            '≥': ' lớn hơn hoặc bằng',
            '≪': ' nhỏ hơn rất nhiều',
            '≫': ' lớn hơn rất nhiều',
            '≡': ' đồng dư',
            '≅': ' đồng dạng',
            '∝': ' tỷ lệ thuận',
            '∑': ' tổng',
            '∏': ' tích',
            '∫': ' tích phân',
            '∬': ' tích phân kép',
            '∭': ' tích phân ba',
            '∮': ' tích phân đường',
            '∯': ' tích phân mặt',
            '∰': ' tích phân thể tích',
            '∇': ' nabla',
            '∂': ' đạo hàm riêng',
            '∆': ' delta',
            '∅': ' tập rỗng',
            '∈': ' thuộc',
            '∉': ' không thuộc',
            '∋': ' chứa',
            '∌': ' không chứa',
            '⊂': ' tập con',
            '⊃': ' tập cha',
            '⊆': ' tập con hoặc bằng',
            '⊇': ' tập cha hoặc bằng',
            '∪': ' hợp',
            '∩': ' giao',
            '∖': ' hiệu',
            '⊕': ' tổng trực tiếp',
            '⊗': ' tích tensor',
            '⊥': ' vuông góc',
            '∥': ' song song',
            '∠': ' góc',
            '∡': ' góc đo',
            '∢': ' góc phẳng',
            '°': ' độ',
            '′': ' phút',
            '″': ' giây',
            '‰': ' phần nghìn',
            '‱': ' phần vạn',

            # Mũi tên
            '→': ' mũi tên phải',
            '←': ' mũi tên trái',
            '↑': ' mũi tên lên',
            '↓': ' mũi tên xuống',
            '↔': ' mũi tên hai chiều',
            '↕': ' mũi tên lên xuống',
            '⇒': ' suy ra',
            '⇐': ' ngược lại',
            '⇔': ' tương đương',
            '⇎': ' không tương đương',

            # Logic
            '∀': ' với mọi',
            '∃': ' tồn tại',
            '∄': ' không tồn tại',
            '∴': ' do đó',
            '∵': ' vì',
            '∧': ' và',
            '∨': ' hoặc',
            '¬': ' không',
            '⊤': ' đúng',
            '⊥': ' sai',

            # Tập số
            'ℕ': ' tập số tự nhiên',
            'ℤ': ' tập số nguyên',
            'ℚ': ' tập số hữu tỷ',
            'ℝ': ' tập số thực',
            'ℂ': ' tập số phức',
            'ℙ': ' tập số nguyên tố',

            # Ký tự khác
            'ℵ': ' aleph',
            'ℶ': ' beth',
            'ℷ': ' gimel',
            'ℸ': ' daleth',
            'ℏ': ' h bar',
            'ℯ': ' e',
            'ℊ': ' g',
            'ℴ': ' o',
            'ℵ': ' aleph',
            '=': ' bằng',
        }

        # Mẫu regex để nhận diện các công thức toán học đơn giản
        self.math_patterns = [
            # Phân số: a/b
            (r'(\d+)/(\d+)', r'\1 chia \2'),
            # Căn bậc hai: √x (cần xử lý trước khi thay thế √)
            (r'√(\w+)', r'căn bậc hai của \1'),
            # Căn bậc n: n√x (cần xử lý trước khi thay thế √)
            (r'(\d+)√(\w+)', r'\1 căn bậc \1 của \2'),
            # Lũy thừa: x^n
            (r'(\w+)\^(\d+)', r'\1 mũ \2'),
            # Tích phân: ∫f(x)dx
            (r'∫([^d]+)d([a-z])', r'tích phân của \1 theo \2'),
            # Đạo hàm: d/dx
            (r'd/(d[a-z])', r'đạo hàm theo \1'),
            # Tổng: Σ
            (r'Σ([^=]+)=([^=]+)', r'tổng của \1 từ \2'),
            # Tích: Π
            (r'Π([^=]+)=([^=]+)', r'tích của \1 từ \2'),
        ]
        # Regex toán tử ASCII an toàn (chỉ áp dụng trong 'math line')
        self.math_ascii_patterns = [
            # tỉ số a:b -> a chia b (giữa các số)
            (re.compile(r'(\b\d+)\s*:\s*(\d+\b)'), r'\1 chia \2'),

            # cộng/trừ/nhân/chia/bằng khi có toán hạng hai bên (chữ/số/ngoặc)
            (re.compile(r'(\b[\w\)\]])\s*\+\s*([\w\(\[]\b)'), r'\1 cộng \2'),
            (re.compile(r'(\b[\w\)\]])\s*-\s*([\w\(\[]\b)'), r'\1 trừ \2'),
            (re.compile(r'(\b[\w\)\]])\s*\*\s*([\w\(\[]\b)'), r'\1 nhân \2'),
            (re.compile(r'(\b[\w\)\]])\s*/\s*([\w\(\[]\b)'), r'\1 chia \2'),
            (re.compile(r'(\b[\w\)\]])\s*=\s*([\w\(\[]\b)'), r'\1 bằng \2'),

            # số âm ở đầu token: -3 -> âm 3 (không đụng gạch nối trong từ)
            (re.compile(r'(?<![\w\.])-([0-9]+)\b'), r'âm \1'),

            # lũy thừa ASCII ^ (giữa biến/số và số/chữ)
            (re.compile(r'([A-Za-z0-9])\s*\^\s*([0-9]+)'), r'\1 mũ \2'),
            (re.compile(r'([A-Za-z0-9])\s*\^\s*([A-Za-z]+)'), r'\1 mũ \2'),

            # 'x' là dấu nhân **chỉ khi** có khoảng trắng hai bên hoặc kẹp giữa số/biến và số/biến
            (re.compile(r'(?<=\b[\da-wyzA-WYZ])\s+x\s+(?=[\da-wyzA-WYZ]\b)'), r' nhân '),
        ]

    # --- Helpers for better math reading ---

    _MATH_LET = r"[A-Za-zα-ωΑ-Ω]"  # chữ Latin + Greek 1 ký tự

    def _normalize_superscripts(self, s: str) -> str:
        """Chuẩn hóa số mũ unicode -> dạng ^n để xử lý/thay thế ổn định hơn."""
        trans = str.maketrans("⁰¹²³⁴⁵⁶⁷⁸⁹", "0123456789")
        s = s.replace("²", "^2").replace("³", "^3").translate(trans)
        return s

    def _insert_multiplication_reading(self, s: str) -> str:
        """
        Chèn ' nhân ' trong các trường hợp thực sự là phép nhân:
        - 2x, 10x  -> 2 nhân x
        - ax, b x  -> a nhân x, b nhân x  (chỉ khi 'a' là biến đơn lẻ, không đứng trong từ)
        - πr, π d  -> π nhân r (hoặc d)
        Chỉ áp dụng khi dòng có 'tính toán học' để tránh false positive trong văn bản thường.
        """
        if not self._is_math_line(s):
            return s

        # 1) số + x  (2x, 10x)
        s = re.sub(r'(?<!\w)(\d+)\s*x\b', r'\1 nhân x', s, flags=re.IGNORECASE)

        # 2) biến đơn + x  (ax, a x) – biến là 1 chữ, và đứng độc lập (không có chữ/số ngay trước)
        s = re.sub(r'(?<![A-Za-z0-9_])([A-Za-z])\s*x\b', r'\1 nhân x', s)         # a x
        s = re.sub(r'(?<![A-Za-z0-9_])([A-Za-z])x\b',     r'\1 nhân x', s)         # ax

        # 3) π + biến đơn (πr, π r)
        s = re.sub(r'π\s*([A-Za-z])\b', r'π nhân \1', s)

        # 4) (số|biến) + x^n  (ax^2, 2x^3) – giữ cụm x^n lại để verbalize sau
        s = re.sub(r'(?<!\w)([A-Za-z]|\d+)\s*x\s*\^\s*(\d+)\b', r'\1 nhân x^\2', s)

        return s


    def _verbalize_exponents(self, s: str) -> str:
        """x^2 -> x mũ 2 (đặt sau khi đã chèn 'nhân' để giữ cụm 'x^2')."""
        s = re.sub(rf"\b({self._MATH_LET})\s*\^\s*(\d+)\b", r"\1 mũ \2", s)
        return s

    def process_special_characters(self, text: str) -> str:
        if not text:
            return text

        s = text

        # (0) Chuẩn hoá superscript -> ^n
        s = self._normalize_superscripts(s)

        # (1) Quy tắc regex “công thức” sẵn có (phân số, căn, tích phân, …)
        for pattern, replacement in self.math_patterns:
            s = re.sub(pattern, replacement, s)

        # (1.5) Chèn 'nhân' giữa số–biến, biến–biến… để tránh đọc dính
        s = self._insert_multiplication_reading(s)

        # (1.6) Verbalize mũ từ dạng ^n
        s = self._verbalize_exponents(s)

        # (2) Nếu là dòng toán, áp dụng thêm ASCII math patterns ( + - * / = ... )
        if self._is_math_line(s):
            for rx, repl in self.math_ascii_patterns:
                s = rx.sub(repl, s)

        # (3) Bản đồ ký hiệu Unicode + fallback ASCII (tránh '.')
        ascii_fallbacks = {
            '+': ' cộng', '-': ' trừ', '=': ' bằng',
            '/': ' chia', '*': ' nhân', '^': ' mũ ',
        }
        merged_map = dict(ascii_fallbacks)
        merged_map.update(self.special_char_map)
        for ch, rep in merged_map.items():
            if ch == '.':
                continue
            s = s.replace(ch, rep)

        # (4) Unicode khác
        s = self._process_unicode_chars(s)

        # (5) Làm sạch
        s = self._clean_text(s)
        return s


    def debug_process(self, text: str) -> Dict[str, str]:
        """
        Debug từng bước xử lý để kiểm tra
        """
        debug_info = {
            'original': text,
            'after_regex': text,
            'after_special_chars': text,
            'after_unicode': text,
            'after_cleaning': text,
            'final': text
        }

        # Bước 1: Regex
        processed_text = text
        for pattern, replacement in self.math_patterns:
            processed_text = re.sub(pattern, replacement, processed_text)
        debug_info['after_regex'] = processed_text

        # Bước 2: Ký tự đặc biệt
        for special_char, replacement in self.special_char_map.items():
            processed_text = processed_text.replace(special_char, replacement)
        debug_info['after_special_chars'] = processed_text

        # Bước 3: Unicode
        processed_text = self._process_unicode_chars(processed_text)
        debug_info['after_unicode'] = processed_text

        # Bước 4: Làm sạch
        processed_text = self._clean_text(processed_text)
        debug_info['after_cleaning'] = processed_text

        debug_info['final'] = processed_text
        return debug_info

    def _process_unicode_chars(self, text: str) -> str:
        """
        Xử lý các ký tự Unicode khác
        """
        processed_chars = []

        for char in text:
            if ord(char) > 127:  # Ký tự không phải ASCII
                try:
                    # Lấy tên Unicode
                    char_name = unicodedata.name(char)

                    # Xử lý các trường hợp đặc biệt
                    if 'SUPERSCRIPT' in char_name:
                        # Số mũ
                        if char_name.endswith('TWO'):
                            processed_chars.append(' mũ hai')
                        elif char_name.endswith('THREE'):
                            processed_chars.append(' mũ ba')
                        elif char_name.endswith('ONE'):
                            processed_chars.append(' mũ một')
                        else:
                            # Lấy số từ tên
                            number = char_name.split()[-1]
                            processed_chars.append(f' mũ {number}')
                    elif 'SUBSCRIPT' in char_name:
                        # Chỉ số dưới
                        if char_name.endswith('TWO'):
                            processed_chars.append(' chỉ số hai')
                        elif char_name.endswith('THREE'):
                            processed_chars.append(' chỉ số ba')
                        elif char_name.endswith('ONE'):
                            processed_chars.append(' chỉ số một')
                        else:
                            number = char_name.split()[-1]
                            processed_chars.append(f' chỉ số {number}')
                    elif 'GREEK' in char_name:
                        # Chữ Hy Lạp
                        greek_name = char_name.split()[-1].lower()
                        processed_chars.append(f' {greek_name}')
                    elif 'MATHEMATICAL' in char_name:
                        # Ký tự toán học
                        math_name = char_name.split()[-1].lower()
                        processed_chars.append(f' {math_name}')
                    else:
                        # Giữ nguyên ký tự nếu không xử lý được
                        processed_chars.append(char)

                except ValueError:
                    # Nếu không lấy được tên Unicode, giữ nguyên
                    processed_chars.append(char)
            else:
                processed_chars.append(char)

        return ''.join(processed_chars)

    def _clean_text(self, text: str) -> str:
        """
        Làm sạch văn bản sau khi xử lý
        """
        # Loại bỏ khoảng trắng thừa
        text = re.sub(r'\s+', ' ', text)

        # Loại bỏ khoảng trắng đầu cuối
        text = text.strip()

        # Xử lý các dấu câu - đảm bảo không có khoảng trắng trước dấu câu
        text = re.sub(r'\s+([.,;:!?])', r'\1', text)

        # Xử lý dấu ngoặc - đảm bảo khoảng trắng phù hợp
        text = re.sub(r'\(\s+', '(', text)
        text = re.sub(r'\s+\)', ')', text)

        # KHÔNG xử lý khoảng trắng giữa các từ đã được xử lý
        # Vì có thể làm hỏng các từ như "mũ hai", "chỉ số ba"

        return text

    def extract_math_objects_from_pptx(self, pptx_file_path: str) -> List[Dict]:
        """
        Trích xuất các đối tượng toán học từ PowerPoint
        """
        try:
            from pptx import Presentation
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import qn

            prs = Presentation(pptx_file_path)
            math_objects = []

            for slide_num, slide in enumerate(prs.slides):
                slide_math = []

                for shape in slide.shapes:
                    # Kiểm tra các đối tượng toán học
                    if hasattr(shape, 'element'):
                        element = shape.element

                        # Tìm các phần tử toán học
                        math_elements = element.findall('.//m:oMath', {'m': 'http://schemas.openxmlformats.org/officeDocument/2006/math'})

                        for math_elem in math_elements:
                            try:
                                # Chuyển đổi XML thành văn bản
                                math_text = self._extract_mathml_text(math_elem)
                                if math_text:
                                    slide_math.append({
                                        'type': 'math_object',
                                        'content': math_text,
                                        'processed_content': self.process_special_characters(math_text)
                                    })
                            except Exception as e:
                                logger.warning(f"Lỗi xử lý phần tử toán học: {e}")

                if slide_math:
                    math_objects.append({
                        'slide_number': slide_num + 1,
                        'math_objects': slide_math
                    })

            return math_objects

        except ImportError:
            logger.error("Không thể import python-pptx. Vui lòng cài đặt: pip install python-pptx")
            return []
        except Exception as e:
            logger.error(f"Lỗi trích xuất đối tượng toán học: {e}")
            return []

    def _extract_mathml_text(self, math_element) -> str:
        """
        Trích xuất văn bản từ phần tử MathML
        """
        try:
            # Chuyển đổi XML thành chuỗi
            xml_str = ET.tostring(math_element, encoding='unicode')

            # Xử lý các thẻ MathML cơ bản
            text_parts = []

            # Tìm tất cả các phần tử văn bản
            for elem in math_element.iter():
                if elem.text and elem.text.strip():
                    text_parts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    text_parts.append(elem.tail.strip())

            return ' '.join(text_parts)

        except Exception as e:
            logger.warning(f"Lỗi trích xuất MathML: {e}")
            return ""

    def process_powerpoint_text(self, pptx_file_path: str) -> Dict:
        """
        Xử lý toàn bộ văn bản từ PowerPoint, bao gồm cả đối tượng toán học
        """
        try:
            from pptx import Presentation

            prs = Presentation(pptx_file_path)
            processed_slides: List[Dict] = []

            # Trích xuất đối tượng toán học
            math_objects = self.extract_math_objects_from_pptx(pptx_file_path)
            math_dict = {obj['slide_number']: obj['math_objects'] for obj in math_objects}

            for i, slide in enumerate(prs.slides):
                slide_num = i + 1

                # Thu thập thông tin shape: loại, tọa độ và văn bản. Chia bảng và
                # các shape khác nhưng vẫn lưu lại vị trí để sắp xếp hợp lý.
                shapes_info: List[Tuple[str, int, int, str]] = []  # (type, top, left, text)
                for shape in slide.shapes:
                    extracted_text = self._extract_text_from_shape(shape)
                    if not extracted_text:
                        continue
                    # loại shape: 'table' hoặc 'text'
                    try:
                        is_table = getattr(shape, 'has_table', False)
                    except Exception:
                        is_table = False
                    shape_type = 'table' if is_table else 'text'
                    top = 0
                    left = 0
                    try:
                        top = int(getattr(shape, 'top', 0))
                        left = int(getattr(shape, 'left', 0))
                    except Exception:
                        pass
                    shapes_info.append((shape_type, top, left, extracted_text))

                # Phân loại lại: tách bảng và các shape khác (non-table)
                table_infos = [info for info in shapes_info if info[0] == 'table']
                non_table_infos = [info for info in shapes_info if info[0] != 'table']

                ordered_texts: List[str] = []
                # Xử lý non-table shapes: nhóm thành cột, xử lý tiêu đề bên trái trước
                if non_table_infos:
                    # Chuyển đổi thành dạng dict để tái sử dụng helpers
                    items = []
                    for t, top, left, text in non_table_infos:
                        items.append({'text': text, 'left': left, 'top': top, 'height': 0})
                    # Tính median width để ước lượng ngưỡng phân cột
                    try:
                        w_med = self._median([s['width'] for s in self._iter_text_shapes_with_pos(slide.shapes) if s.get('width', 0) > 0])
                    except Exception:
                        w_med = 0
                    # Tính col threshold tương tự fallback
                    col_thr = max(int((w_med or 1) * 0.4), int(prs.slide_width / 80))
                    cols = self._group_columns(items, col_thr)
                    cols.sort(key=lambda c: c['x'])
                    if len(cols) > 1:
                        # Xác định cột chứa tiêu đề/nội dung chính.
                        # Nếu chỉ có 2 cột, ưu tiên chọn cột bên trái làm tiêu đề trừ khi cột đó chỉ chứa nhãn số (1–3 ký tự),
                        # khi đó chọn cột còn lại làm tiêu đề. Nếu có nhiều hơn 2 cột, chọn cột có số phần tử ít nhất.
                        header_col_idx = 0
                        if len(cols) == 2:
                            # Kiểm tra xem tất cả phần tử ở cột 0 có phải là nhãn số hay không
                            all_numeric_in_col0 = all(
                                self._is_numeric_badge(item['text']) for item in cols[0]['items']
                            )
                            if all_numeric_in_col0:
                                header_col_idx = 1
                        elif len(cols) > 2:
                            # Tìm cột có số phần tử ít nhất
                            min_count = min(len(c['items']) for c in cols)
                            candidate_cols = [idx for idx, c in enumerate(cols) if len(c['items']) == min_count]
                            header_col_idx = candidate_cols[0] if candidate_cols else 0
                        # Lấy các mục tiêu đề và danh sách
                        left_col_items = cols[header_col_idx]['items']
                        right_items: List[Dict] = []
                        for idx, col in enumerate(cols):
                            if idx != header_col_idx:
                                right_items += col['items']
                        # Sắp xếp theo top cho tiêu đề và danh sách
                        left_texts = [s['text'] for s in sorted(left_col_items, key=lambda s: s['top'])]
                        right_sorted = sorted(right_items, key=lambda s: s['top'])
                        # Nếu danh sách chỉ bao gồm các nhãn số và số nhãn bằng số tiêu đề, ghép cặp chúng
                        numeric_badges = [item for item in right_sorted if self._is_numeric_badge(item['text'])]
                        if numeric_badges and len(numeric_badges) == len(right_sorted) == len(left_texts):
                            ordered_texts = []
                            for badge_item, left_txt in zip(numeric_badges, left_texts):
                                ordered_texts.append(f"{badge_item['text']} {left_txt}")
                        else:
                            # Ghép nhãn số với mô tả liền sau nó
                            right_lines: List[str] = []
                            for item in right_sorted:
                                txt = item['text']
                                if self._is_numeric_badge(txt) and right_lines:
                                    prev = right_lines.pop()
                                    right_lines.append(f"{txt} {prev}")
                                else:
                                    right_lines.append(txt)
                            ordered_texts = left_texts + right_lines
                    else:
                        # Chỉ một cột, sắp xếp theo top
                        ordered_texts = [s['text'] for s in sorted(cols[0]['items'], key=lambda s: s['top'])]
                # Thêm bảng (nếu có) vào cuối danh sách
                for (_, _, _, table_text) in sorted(table_infos, key=lambda x: (x[1], x[2])):
                    ordered_texts.append(table_text)
                slide_text = "\n".join(filter(None, ordered_texts))

                # Xử lý văn bản thông thường
                processed_text = self.process_special_characters(slide_text)

                # Thêm đối tượng toán học nếu có
                if slide_num in math_dict:
                    math_texts: List[str] = []
                    for math_obj in math_dict[slide_num]:
                        math_texts.append(math_obj['processed_content'])
                    if math_texts:
                        processed_text += " " + " ".join(math_texts)

                processed_slides.append({
                    'slide_number': slide_num,
                    'original_text': slide_text,
                    'processed_text': processed_text,
                    'has_math_objects': slide_num in math_dict
                })

            return {
                'slides': processed_slides,
                'total_slides': len(processed_slides),
                'slides_with_math': len([s for s in processed_slides if s['has_math_objects']])
            }

        except Exception as e:
            logger.error(f"Lỗi xử lý PowerPoint: {e}")
            return {
                'slides': [],
                'total_slides': 0,
                'slides_with_math': 0,
                'error': str(e)
            }

    # ---------- Layout helpers (vị trí/khối) ----------
    def _is_numeric_badge(self, txt: str) -> bool:
        return bool(re.match(r"^\s*\d{1,3}[.)-]?\s*$", (txt or "").strip()))

    def _median(self, vals):
        vals = sorted(int(v) for v in vals if v is not None and v > 0)
        if not vals:
            return 0
        n = len(vals)
        return vals[n // 2] if n % 2 else (vals[n // 2 - 1] + vals[n // 2]) // 2

    def _iter_text_shapes_with_pos(self, shapes, dx=0, dy=0):
        for shp in shapes:
            try:
                st = getattr(shp, 'shape_type', None)
            except Exception:
                st = None

            if MSO_SHAPE_TYPE is not None and st == MSO_SHAPE_TYPE.GROUP:
                ox, oy = int(getattr(shp, 'left', 0)), int(getattr(shp, 'top', 0))
                yield from self._iter_text_shapes_with_pos(shp.shapes, dx + ox, dy + oy)
                continue

            # text content
            txt = ''
            if getattr(shp, 'has_text_frame', False) and getattr(shp, 'text_frame', None) is not None:
                try:
                    txt = shp.text_frame.text or ''
                except Exception:
                    txt = ''
            elif getattr(shp, 'text', None):
                txt = shp.text or ''
            txt = (txt or '').strip()
            if not txt:
                continue

            yield {
                'text': txt,
                'left': int(getattr(shp, 'left', 0)) + dx,
                'top': int(getattr(shp, 'top', 0)) + dy,
                'width': int(getattr(shp, 'width', 0)),
                'height': int(getattr(shp, 'height', 0)),
            }

    def _group_columns(self, items, col_thr):
        cols = []
        for s in sorted(items, key=lambda a: a['left']):
            placed = False
            for c in cols:
                if abs(s['left'] - c['x']) <= col_thr:
                    c['items'].append(s)
                    c['x'] = (c['x'] * (len(c['items']) - 1) + s['left']) // len(c['items'])
                    placed = True
                    break
            if not placed:
                cols.append({'x': s['left'], 'items': [s]})
        return cols

    def _group_rows(self, items, row_thr):
        rows = []
        for s in sorted(items, key=lambda a: a['top']):
            cy = s['top'] + max(1, s.get('height', 0)) // 2
            placed = False
            for r in rows:
                if abs(cy - r['cy']) <= row_thr:
                    r['items'].append(s)
                    r['cy'] = (r['cy'] * (len(r['items']) - 1) + cy) // len(r['items'])
                    placed = True
                    break
            if not placed:
                rows.append({'cy': cy, 'items': [s]})
        return rows

    def _compose_row_text(self, row_items):
        row_items = sorted(row_items, key=lambda s: s['left'])
        nums = [s for s in row_items if self._is_numeric_badge(s['text'])]
        others = [s for s in row_items if not self._is_numeric_badge(s['text'])]
        if nums and others:
            badge = nums[0]['text'].strip()
            other_text = ' '.join(s['text'].strip() for s in others if s['text'].strip())
            return (badge + ' ' + other_text).strip() if other_text else badge
        return ' '.join(s['text'].strip() for s in row_items if s['text'].strip())

    def _extract_text_from_shape(self, shape) -> str:
        """
        Trích xuất text từ shape. Hỗ trợ: textbox, bảng, group và placeholder.
        Đệ quy qua group. Ghép paragraphs/runs để không mất ký tự.
        """
        try:
            # Text frame (textbox, placeholder)
            if getattr(shape, 'has_text_frame', False):
                parts: List[str] = []
                tf = shape.text_frame
                try:
                    for p in tf.paragraphs:
                        if getattr(p, 'runs', None):
                            run_text = ''.join([r.text for r in p.runs if getattr(r, 'text', None)])
                            parts.append(run_text if run_text else p.text)
                        else:
                            parts.append(p.text)
                except Exception:
                    if getattr(shape, 'text', None):
                        parts.append(shape.text)
                return self._clean_text(' '.join([t for t in parts if t]))

            # Table - đọc từ trái qua phải, từ trên xuống dưới
            if getattr(shape, 'has_table', False):
                table_texts: List[str] = []
                for row in shape.table.rows:
                    row_texts: List[str] = []
                    for cell in row.cells:
                        cell_text = ""
                        if getattr(cell, 'text_frame', None):
                            cell_parts: List[str] = []
                            for p in cell.text_frame.paragraphs:
                                if getattr(p, 'runs', None):
                                    run_text = ''.join([r.text for r in p.runs if getattr(r, 'text', None)])
                                    cell_parts.append(run_text if run_text else p.text)
                                else:
                                    cell_parts.append(p.text)
                            cell_text = ' '.join(cell_parts) if cell_parts else ""
                        elif getattr(cell, 'text', None):
                            cell_text = cell.text
                        if cell_text.strip():
                            row_texts.append(cell_text.strip())
                    if row_texts:
                        table_texts.append(' '.join(row_texts))
                return '\n'.join(table_texts)

            # Group: traverse children
            if getattr(shape, 'shapes', None) is not None and getattr(shape, 'shape_type', None) is not None and 'GROUP' in str(shape.shape_type):
                sub_texts: List[str] = []
                for sub in shape.shapes:
                    t = self._extract_text_from_shape(sub)
                    if t:
                        sub_texts.append(t)
                return self._clean_text(' '.join(sub_texts))

            # Fallback
            if getattr(shape, 'text', None):
                return self._clean_text(shape.text)
        except Exception as e:
            logger.debug(f"_extract_text_from_shape error: {e}")
        return ''


# Hàm tiện ích để sử dụng nhanh
def process_math_text(text: str) -> str:
    """
    Hàm tiện ích để xử lý nhanh văn bản chứa công thức toán học
    """
    processor = MathFormulaProcessor()
    return processor.process_special_characters(text)


def process_powerpoint_file(pptx_file_path: str) -> Dict:
    """
    Hàm tiện ích để xử lý toàn bộ file PowerPoint
    """
    processor = MathFormulaProcessor()
    return processor.process_powerpoint_text(pptx_file_path)