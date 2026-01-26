# app/services/pptx_service.py
import os
import time
import tempfile
import subprocess
from shutil import which
from typing import Any

from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image

try:
    import pytesseract
except Exception:
    pytesseract = None

from src.utils.math_formula_processor import MathFormulaProcessor, process_math_text
from app.config import get_config


def _as_path(p: Any) -> str | None:
    """
    Chấp nhận:
    - đường dẫn str / PathLike
    - object có .name (ví dụ werkzeug FileStorage đã save ra file, hoặc object tương tự)
    """
    if p is None:
        return None
    if isinstance(p, (str, os.PathLike)):
        return str(p)
    return getattr(p, "name", None)


def convert_pptx_to_images(pptx_path: str, dpi: int = 220) -> list[str]:
    """
    PPTX -> PDF bằng LibreOffice -> PNG bằng pdf2image.
    Trả về list đường dẫn ảnh slide PNG.
    """
    cfg = get_config()
    lo = cfg.LIBREOFFICE_PATH

    if not lo or which(lo) is None:
        raise RuntimeError(
            f"Không tìm thấy LibreOffice tại: {lo}. "
            f"Hãy set biến môi trường LIBREOFFICE_PATH hoặc cài LibreOffice đúng."
        )

    tmpdir = tempfile.mkdtemp(prefix="pptx2img_", dir=cfg.TMP_DIR if os.path.isdir(cfg.TMP_DIR) else None)

    subprocess.run(
        [lo, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )

    pdf_path = os.path.join(tmpdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")

    # Chờ LO flush file
    time.sleep(0.7)

    if not os.path.exists(pdf_path):
        raise RuntimeError(f"Convert PPTX->PDF thất bại, không thấy file PDF: {pdf_path}")

    poppler_path = cfg.POPPLER_PATH
    kwargs = dict(dpi=dpi, output_folder=tmpdir, fmt="png")
    if poppler_path:
        kwargs["poppler_path"] = poppler_path

    images = convert_from_path(pdf_path, **kwargs)

    out: list[str] = []
    for i, img in enumerate(images, 1):
        p = os.path.join(tmpdir, f"slide-{i:02d}.png")
        img.save(p)
        out.append(p)
    return out


def extract_slides_from_pptx(pptx_file_or_path: Any, dpi: int = 220) -> list[dict]:
    """
    Output chuẩn hoá:
    [
      {
        "slide_number": int,
        "text": str,
        "image_path": str | None,
        "has_math_objects": bool
      }, ...
    ]
    """
    pptx_path = _as_path(pptx_file_or_path)
    if not pptx_path or not os.path.exists(pptx_path):
        raise RuntimeError("Không tìm thấy file PowerPoint hợp lệ.")

    imgs = convert_pptx_to_images(pptx_path, dpi=dpi)
    slides: list[dict] = []

    # Ưu tiên MathFormulaProcessor
    mp = MathFormulaProcessor()
    res = mp.process_powerpoint_text(pptx_path)

    if not res.get("error"):
        for s in res["slides"]:
            idx = s["slide_number"] - 1
            slides.append(
                {
                    "slide_number": s["slide_number"],
                    "text": s["processed_text"],
                    "image_path": imgs[idx] if 0 <= idx < len(imgs) else None,
                    "has_math_objects": bool(s.get("has_math_objects")),
                }
            )
        return slides

    # Fallback: đọc text thô + OCR
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        chunks: list[str] = []
        for shp in slide.shapes:
            try:
                if getattr(shp, "has_text_frame", False) and shp.text_frame:
                    t = (shp.text_frame.text or "").strip()
                    if t:
                        chunks.append(t)
            except Exception:
                pass

        text = process_math_text("\n".join(chunks).strip())

        if (not text) and (i < len(imgs)) and (pytesseract is not None):
            try:
                ocr = pytesseract.image_to_string(Image.open(imgs[i]), lang="vie+eng")
                text = process_math_text((ocr or "").strip())
            except Exception:
                text = ""

        slides.append(
            {
                "slide_number": i + 1,
                "text": text,
                "image_path": imgs[i] if i < len(imgs) else None,
                "has_math_objects": False,
            }
        )
    return slides


def format_slides_as_text(slides: list[dict]) -> str:
    """
    Giống _format_slides_as_text bên Gradio để đổ vào editor.
    """
    lines: list[str] = []
    for s in slides:
        n = s.get("slide_number")
        lines.append(f"## Slide {n}")
        lines.append((s.get("text") or "").strip())
        lines.append("")
    return "\n".join(lines).strip()
